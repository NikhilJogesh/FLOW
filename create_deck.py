import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

# Colors
BG_COLOR = hex_to_rgb('#0B0F19')
TEXT_PRIMARY = hex_to_rgb('#FFFFFF')
TEXT_SECONDARY = hex_to_rgb('#94A3B8')
BLUE_ELECTRIC = hex_to_rgb('#3B82F6')
AMBER_WARNING = hex_to_rgb('#F59E0B')
RED_FAILURE = hex_to_rgb('#EF4444')
EMERALD_SUCCESS = hex_to_rgb('#10B981')

prs = Presentation()
# Set widescreen 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text, color=TEXT_PRIMARY, font_size=44, top=0.5):
    txBox = slide.shapes.add_textbox(Inches(1), Inches(top), Inches(11.333), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    return txBox

def add_subtitle(slide, text, top=1.5, color=TEXT_SECONDARY, size=24, bold=False):
    txBox = slide.shapes.add_textbox(Inches(1), Inches(top), Inches(11.333), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return txBox

def add_notes(slide, notes):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes

# ----------------------------------------------------
# SLIDE 1 — TITLE
# ----------------------------------------------------
slide1 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide1)
add_title(slide1, "FLOW", color=BLUE_ELECTRIC, font_size=80, top=2.5)
add_subtitle(slide1, "Predictive Transit Recovery", top=3.8, color=TEXT_PRIMARY, size=32, bold=True)
add_subtitle(slide1, "\"Don't reroute after failure.\nRecover before it.\"", top=4.5, color=TEXT_SECONDARY, size=24)
add_notes(slide1, "Speaker Notes:\nWelcome everyone. We are here to introduce FLOW, our predictive transit recovery engine. Our core philosophy is simple: don't wait for the connection to fail, predict the failure and recover before it happens.")

# ----------------------------------------------------
# SLIDE 2 — THE PROBLEM
# ----------------------------------------------------
slide2 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide2)
add_title(slide2, "The fastest route can still fail.")

cascade = [
    ("DELAY", RED_FAILURE),
    ("MISSED TRANSFER", RED_FAILURE),
    ("STRANDED COMMUTER", TEXT_SECONDARY),
    ("PANIC REROUTING", AMBER_WARNING),
    ("EVERYONE TAKES THE SAME ALTERNATIVE", AMBER_WARNING),
    ("SECONDARY BOTTLENECK", RED_FAILURE)
]

for i, (text, color) in enumerate(cascade):
    y = 2.0 + (i * 0.8)
    # Add text
    tb = add_subtitle(slide2, text, top=y, color=color, size=24, bold=True)
    tb.left = Inches(4.5)
    # Add arrow if not last
    if i < len(cascade) - 1:
        arr = slide2.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.6), Inches(y + 0.5), Inches(0.2), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = TEXT_SECONDARY
        arr.line.fill.background()

add_notes(slide2, "Speaker Notes:\nConventional routing asks 'what is the fastest route'. But what happens when there's a delay? The commuter misses their transfer, and the system panic reroutes them. Because every app uses the same fastest-path algorithm, everyone is routed to the exact same alternative, creating a massive secondary bottleneck.")

# ----------------------------------------------------
# SLIDE 3 — THE FLOW DIFFERENCE
# ----------------------------------------------------
slide3 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide3)
add_title(slide3, "FLOW doesn't ask what's fastest.")

# Traditional
txBoxT = slide3.shapes.add_textbox(Inches(1), Inches(2.5), Inches(5), Inches(2))
tfT = txBoxT.text_frame
pT = tfT.paragraphs[0]
pT.text = "TRADITIONAL\n\n\"What is the fastest route?\""
pT.font.name = 'Arial'
pT.font.size = Pt(28)
pT.font.color.rgb = TEXT_SECONDARY
pT.alignment = PP_ALIGN.CENTER

# FLOW
txBoxF = slide3.shapes.add_textbox(Inches(7.333), Inches(2.5), Inches(5), Inches(2))
tfF = txBoxF.text_frame
pF = tfF.paragraphs[0]
pF.text = "FLOW\n\n\"What route is most likely to succeed?\""
pF.font.name = 'Arial'
pF.font.size = Pt(28)
pF.font.color.rgb = BLUE_ELECTRIC
pF.font.bold = True
pF.alignment = PP_ALIGN.CENTER

# PREDICT RECOVER REDISTRIBUTE
txBoxM = slide3.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(1))
tfM = txBoxM.text_frame
pM = tfM.paragraphs[0]
pM.text = "PREDICT  →  RECOVER  →  REDISTRIBUTE"
pM.font.name = 'Arial'
pM.font.size = Pt(36)
pM.font.color.rgb = EMERALD_SUCCESS
pM.font.bold = True
pM.alignment = PP_ALIGN.CENTER

add_notes(slide3, "Speaker Notes:\nFLOW flips the paradigm. Instead of asking what is fastest in a perfect world, we ask: what route has the highest probability of success right now? We do this in three steps: Predict the failure, Recover the commuter with viable options, and Redistribute the network load.")


# ----------------------------------------------------
# SLIDE 4 — CONNECTION CONFIDENCE
# ----------------------------------------------------
slide4 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide4)
add_title(slide4, "FLOW sees the failure coming.")

add_subtitle(slide4, "Connection Buffer = Predicted Arrival − Required Transfer Time", top=1.8, size=20)
add_subtitle(slide4, "Risk Margin = Connection Buffer − Historical P95 Delay", top=2.4, size=20)
add_subtitle(slide4, "Connection Confidence = deterministic normalized score", top=3.0, size=20)

# Visual Example
y_base = 4.5
add_subtitle(slide4, "100% CONFIDENCE", top=y_base, color=EMERALD_SUCCESS, size=24, bold=True)
arr1 = slide4.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2), Inches(y_base+0.5), Inches(0.2), Inches(0.3))
arr1.fill.solid(); arr1.fill.fore_color.rgb = TEXT_SECONDARY
add_subtitle(slide4, "Bus predicted +12 min late", top=y_base+0.8, color=AMBER_WARNING, size=24, bold=True)
arr2 = slide4.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2), Inches(y_base+1.3), Inches(0.2), Inches(0.3))
arr2.fill.solid(); arr2.fill.fore_color.rgb = TEXT_SECONDARY
add_subtitle(slide4, "60% CONFIDENCE → AT RISK", top=y_base+1.6, color=RED_FAILURE, size=28, bold=True)

add_notes(slide4, "Speaker Notes:\nHow do we predict it? Through Connection Confidence. We take the predicted arrival, subtract the transfer time, and penalize by the route's historical P95 delay. This gives us a deterministic risk margin. In our MVP, if a bus is delayed 12 minutes, confidence immediately drops to 60%, triggering an At-Risk state before the commuter is stranded.")


# ----------------------------------------------------
# SLIDE 5 — THE COMMUTER EXPERIENCE
# ----------------------------------------------------
slide5 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide5)
add_title(slide5, "From panic to a recovery decision.")

flow = [
    ("PLAN", BLUE_ELECTRIC),
    ("96/100 CONFIDENCE", EMERALD_SUCCESS),
    ("DELAY DETECTED", AMBER_WARNING),
    ("CONFIDENCE DROPS", RED_FAILURE),
    ("EXPLANATION", TEXT_SECONDARY),
    ("FALLBACK OPTIONS", BLUE_ELECTRIC),
    ("USER SELECTS", EMERALD_SUCCESS),
    ("RECOVERED", EMERALD_SUCCESS)
]

for i, (text, color) in enumerate(flow):
    x = 1.0 + (i * 1.5)
    tb = slide5.shapes.add_textbox(Inches(x), Inches(3.5), Inches(1.4), Inches(1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if i < len(flow) - 1:
        arr = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+1.15), Inches(3.8), Inches(0.25), Inches(0.15))
        arr.fill.solid(); arr.fill.fore_color.rgb = TEXT_SECONDARY; arr.line.fill.background()

add_subtitle(slide5, "Explainability: WHAT CHANGED | WHY IT MATTERS | JOURNEY IMPACT | RECOMMENDATION", top=5.5, size=20, color=TEXT_SECONDARY)

add_notes(slide5, "Speaker Notes:\nThe commuter experience changes from a reactive panic to a proactive decision. The user is actively informed with an explainability payload: What changed, why it matters, and clear, scored fallback options. They select a fallback and recover their journey instantly.")


# ----------------------------------------------------
# SLIDE 6 — LIVE DEMO SCENARIO
# ----------------------------------------------------
slide6 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide6)
add_title(slide6, "Aarav's connection is about to fail.")

timeline = [
    ("08:17", "Delay detected"),
    ("08:17", "FLOW detects risk"),
    ("08:17", "Confidence → 60%"),
    ("08:17", "Fallback options generated"),
    ("08:18", "Aarav switches route"),
    ("08:38", "Baseline connection becomes impossible")
]

for i, (time, desc) in enumerate(timeline):
    y = 2.0 + (i * 0.7)
    tb_time = add_subtitle(slide6, time, top=y, color=BLUE_ELECTRIC, size=20, bold=True)
    tb_time.left = Inches(1)
    
    color_desc = RED_FAILURE if i == 5 else TEXT_PRIMARY
    tb_desc = add_subtitle(slide6, desc, top=y, color=color_desc, size=20)
    tb_desc.left = Inches(2.5)

add_subtitle(slide6, "Prediction Lead Time = 21 minutes", top=6.2, color=EMERALD_SUCCESS, size=28, bold=True)
add_subtitle(slide6, "SIMULATED RESULT", top=6.8, color=TEXT_SECONDARY, size=12, bold=True)

add_notes(slide6, "Speaker Notes:\n(Transition to Live Demo of P0 Flow after this slide).\nIn our demo scenario, FLOW gives Aarav a 21-minute prediction lead time. He recovers his journey at 8:18 AM. Without FLOW, he wouldn't realize his connection was impossible until 8:38 AM.")


# ----------------------------------------------------
# SLIDE 7 — WHY ONE USER ISN'T ENOUGH
# ----------------------------------------------------
slide7 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide7)
add_title(slide7, "Now scale the problem.")

add_subtitle(slide7, "1 commuter  →  1,000 commuters  →  5,000 commuters  →  City Network", top=2.0, color=BLUE_ELECTRIC, size=24, bold=True)

text7 = "If everyone chooses the fastest fallback,\nthe fastest alternative becomes the next bottleneck.\n\nFLOW instead considers:\n\nUSER PREFERENCES\n+ ROUTE CAPACITY\n+ RELIABILITY\n+ TIME\n+ COST\n+ ECO IMPACT"
add_subtitle(slide7, text7, top=3.2, color=TEXT_PRIMARY, size=22)

add_notes(slide7, "Speaker Notes:\nSaving one user is easy. But at city scale, if every user takes the single fastest fallback route, we just create a new secondary bottleneck. FLOW solves this by orchestrating demand. We rank fallbacks based on individualized commuter weights and real-time route capacity constraints.")


# ----------------------------------------------------
# SLIDE 8 — REDISTRIBUTION
# ----------------------------------------------------
slide8 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide8)
add_title(slide8, "FLOW doesn't just reroute. It redistributes.")

# Baseline Side
tb_base = add_subtitle(slide8, "BASELINE (1,000 commuters)", top=2.0, color=TEXT_SECONDARY, size=20, bold=True)
add_subtitle(slide8, "Route B → capacity 200\n→ 500% utilization\n→ 800 failures\n→ 20% success", top=2.6, color=RED_FAILURE, size=24, bold=True)

# FLOW Side
tb_flow = add_subtitle(slide8, "FLOW", top=2.0, color=BLUE_ELECTRIC, size=20, bold=True)
tb_flow.left = Inches(7)
tb_flow2 = add_subtitle(slide8, "Route B → 200 / 200\nRoute C → 800 / 800\nRoute D → 0 / 350\n\n→ 1,000 successful\n→ 100% success", top=2.6, color=EMERALD_SUCCESS, size=24, bold=True)
tb_flow2.left = Inches(7)

add_subtitle(slide8, "SIMULATED RESULT", top=6.8, color=TEXT_SECONDARY, size=12, bold=True)

add_notes(slide8, "Speaker Notes:\nIn our deterministic 1,000 commuter simulation, the baseline routing forces 1,000 people onto a 200-capacity route, resulting in 800 failures. FLOW redistributes them perfectly across Route B and Route C, yielding 100% success. (Transition to Live Demo of Dashboard).")


# ----------------------------------------------------
# SLIDE 9 — COMPLEX SYNTHETIC CITY
# ----------------------------------------------------
slide9 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide9)
add_title(slide9, "From one journey to a network.")

features = [
    "Scalable synthetic network (50–500+ nodes)",
    "Multimodal routes with transfer delays",
    "Generated Hubs & Interchanges",
    "Alternative Pathfinding",
    "Route capacities & bottlenecks",
    "Deterministic PRNG simulation",
    "Disruption modeling"
]
y_val = 2.0
for feat in features:
    add_subtitle(slide9, "• " + feat, top=y_val, color=TEXT_PRIMARY, size=22)
    y_val += 0.5

# Right side visual representation placeholder
box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(5), Inches(4.5))
box.fill.solid(); box.fill.fore_color.rgb = hex_to_rgb('#1E293B')
box.line.color.rgb = BLUE_ELECTRIC
tf_box = box.text_frame
tf_box.text = "[ Synthetic Network Visualization ]\n\nNodes, Edges, Hubs dynamically generated and routed."
tf_box.paragraphs[0].font.color.rgb = BLUE_ELECTRIC
tf_box.paragraphs[0].alignment = PP_ALIGN.CENTER

add_notes(slide9, "Speaker Notes:\nTo prove this works beyond a 3-route toy problem, we built a fully dynamic, scalable synthetic transit network generator. We can spin up a 500-node, multimodal transit city in milliseconds and run our predictive routing and simulation on top of it, identifying alternative paths automatically.")


# ----------------------------------------------------
# SLIDE 10 — TECHNICAL ARCHITECTURE
# ----------------------------------------------------
slide10 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide10)
add_title(slide10, "Under the hood.")

arch = [
    "Synthetic Transit Network",
    "Graph / Routing Engine",
    "Connection Confidence",
    "Journey State Machine",
    "Fallback Scoring",
    "Capacity-Aware Distribution",
    "Commuter UI + Operator Dashboard"
]
for i, step in enumerate(arch):
    y = 1.8 + (i * 0.7)
    add_subtitle(slide10, step, top=y, color=BLUE_ELECTRIC, size=20, bold=True)
    if i < len(arch) - 1:
        arr = slide10.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2), Inches(y+0.45), Inches(0.15), Inches(0.2))
        arr.fill.solid(); arr.fill.fore_color.rgb = TEXT_SECONDARY; arr.line.fill.background()

tech_stack = "Stack:\n• Next.js\n• TypeScript\n• Prisma\n• SQLite\n• Deterministic PRNG simulation\n• React SVG network visualization"
tb_tech = add_subtitle(slide10, tech_stack, top=2.5, color=TEXT_PRIMARY, size=20)
tb_tech.left = Inches(7)

add_notes(slide10, "Speaker Notes:\nOur architecture is a modular Next.js monolith using TypeScript, SQLite, and Prisma. We built a custom graph engine and fallback scoring heuristic entirely from scratch, utilizing a deterministic PRNG so that our simulations are reproducible.")


# ----------------------------------------------------
# SLIDE 11 — WHAT WE ACTUALLY BUILT
# ----------------------------------------------------
slide11 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide11)
add_title(slide11, "Not a mockup.")

built_left = [
    "✓ Connection Confidence Engine",
    "✓ Predictive Delay Detection",
    "✓ Journey State Machine",
    "✓ Explainability Engine",
    "✓ Weighted Fallback Ranking",
    "✓ Proactive Recovery UI"
]

built_right = [
    "✓ Multiple Disruption Recovery",
    "✓ Baseline Comparison",
    "✓ 1,000+ Commuter Simulation",
    "✓ Capacity-Aware Redistribution",
    "✓ Dynamic Synthetic Network",
    "✓ Operator Dashboard"
]

for i, item in enumerate(built_left):
    add_subtitle(slide11, item, top=2.0 + (i * 0.6), color=EMERALD_SUCCESS, size=22, bold=True)
    
for i, item in enumerate(built_right):
    tb = add_subtitle(slide11, item, top=2.0 + (i * 0.6), color=EMERALD_SUCCESS, size=22, bold=True)
    tb.left = Inches(7)

add_notes(slide11, "Speaker Notes:\nEverything you saw today is fully implemented. We didn't use mockups. From the predictive engine to the capacity-aware network redistribution, we built a functional prototype of a completely new way to handle transit routing.")


# ----------------------------------------------------
# SLIDE 12 — THE VISION / CLOSE
# ----------------------------------------------------
slide12 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide12)
add_title(slide12, "Transit shouldn't just recover. It should recover intelligently.", font_size=40)

txBoxEnd = slide12.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(3))
tfEnd = txBoxEnd.text_frame
pEnd1 = tfEnd.paragraphs[0]
pEnd1.text = "PREDICT  →  RECOVER  →  REDISTRIBUTE\n\n"
pEnd1.font.name = 'Arial'
pEnd1.font.size = Pt(36)
pEnd1.font.color.rgb = BLUE_ELECTRIC
pEnd1.font.bold = True
pEnd1.alignment = PP_ALIGN.CENTER

pEnd2 = tfEnd.add_paragraph()
pEnd2.text = "FLOW turns transit routing from a shortest-path problem into a journey-success problem.\n\n"
pEnd2.font.name = 'Arial'
pEnd2.font.size = Pt(24)
pEnd2.font.color.rgb = TEXT_PRIMARY
pEnd2.alignment = PP_ALIGN.CENTER

pEnd3 = tfEnd.add_paragraph()
pEnd3.text = "\"Don't wait for the connection to fail.\""
pEnd3.font.name = 'Arial'
pEnd3.font.size = Pt(32)
pEnd3.font.color.rgb = EMERALD_SUCCESS
pEnd3.font.bold = True
pEnd3.alignment = PP_ALIGN.CENTER

add_notes(slide12, "Speaker Notes:\nBy shifting from fastest-path to journey-success, FLOW orchestrates demand before bottlenecks happen. Don't wait for the connection to fail. Thank you.")

prs.save("FLOW_Predictive_Transit_Recovery.pptx")
